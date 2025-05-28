import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def natural_sort_key(s):
    """按数字顺序排序文件名"""
    return [int(text) if text.isdigit() else text.lower() 
            for text in re.split(r'(\d+)', s)]

def resize_images_and_create_pptx(input_dir):
    """
    将指定目录中的所有JPG图片添加到PPTX演示文稿
    按照文件名数字顺序排序（如1.jpg, 2.jpg, 10.jpg等）
    
    参数:
        input_dir: 输入图片目录
    """
    # 收集所有JPG文件
    jpg_files = []
    for filename in os.listdir(input_dir):
        if filename.lower().endswith(('.jpg', '.jpeg')):
            jpg_files.append(filename)
    
    if not jpg_files:
        print("目录中没有找到JPG图片文件")
        return
    
    # 按照数字顺序排序文件名
    jpg_files.sort(key=natural_sort_key)
    
    # 创建PPTX演示文稿
    pptx_path = os.path.join(input_dir, "图片演示文稿.pptx")
    prs = Presentation()
    
    # 设置16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 为每张图片创建一张幻灯片
    for filename in jpg_files:
        img_path = os.path.join(input_dir, filename)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加图片到幻灯片，自动适应
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        print(f"已添加: {filename}")
    
    # 保存PPTX文件
    prs.save(pptx_path)
    print(f"完成！PPTX文件已保存: {pptx_path}")

if __name__ == "__main__":
    input_dir = input("请输入包含JPG图片的目录路径: ")
    resize_images_and_create_pptx(input_dir)
